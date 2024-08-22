from pptx import Presentation
# Open a new presentation object
prs = Presentation()
# Create the first page object of the presentation file
title_slide_layout = prs.slide_layouts[0]
# Use the method in the presentation object to put the first page object created in the previous line into the presentation
slide = prs.slides.add_slide(title_slide_layout)
# Set the title of the first page
title = slide.shapes.title
title.text = "Hello python pptx~"
# Set the subtitle of the first page
subtitle = slide.placeholders[1]
subtitle.text = "Quite Cool"
# Save the presentation object
prs.save("pythonppt.pptx")